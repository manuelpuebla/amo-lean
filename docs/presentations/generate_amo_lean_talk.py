#!/usr/bin/env python3
"""
AMO-Lean Technical Presentation Generator
Generates a PowerPoint presentation for a technical talk about AMO-Lean.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Professional dark theme
COLORS = {
    'bg_dark': RgbColor(0x1a, 0x1a, 0x2e),      # Dark blue-black
    'bg_medium': RgbColor(0x16, 0x21, 0x3e),    # Medium dark blue
    'accent_blue': RgbColor(0x00, 0xd4, 0xff),  # Cyan accent
    'accent_green': RgbColor(0x00, 0xff, 0x88), # Green accent
    'accent_orange': RgbColor(0xff, 0x88, 0x00),# Orange accent
    'accent_red': RgbColor(0xff, 0x44, 0x44),   # Red accent
    'accent_purple': RgbColor(0xaa, 0x88, 0xff),# Purple accent
    'text_white': RgbColor(0xff, 0xff, 0xff),   # White text
    'text_gray': RgbColor(0xaa, 0xaa, 0xaa),    # Gray text
    'code_bg': RgbColor(0x0d, 0x11, 0x17),      # Code background
}

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent_blue']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, number=None):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_medium'])

    # Number circle if provided
    if number:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.25), Inches(1.5), Inches(1.5), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['accent_blue']
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['bg_dark']
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(15)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, accent_color=None):
    """Add a content slide with bullet points."""
    if accent_color is None:
        accent_color = COLORS['accent_blue']

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['bg_medium']
    title_bar.line.fill.background()

    # Accent line
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(10), Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent_color
    accent_line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle indentation levels
        if bullet.startswith('  - '):
            p.text = bullet[4:]
            p.level = 1
            p.font.size = Pt(18)
        elif bullet.startswith('- '):
            p.text = bullet[2:]
            p.level = 0
            p.font.size = Pt(20)
        else:
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

        p.font.color.rgb = COLORS['text_white']
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    return slide

def add_code_slide(prs, title, code, language="lean"):
    """Add a slide with code."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    # Code box background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.3))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = COLORS['code_bg']
    code_bg.line.color.rgb = COLORS['accent_blue']
    code_bg.line.width = Pt(1)

    # Language tag
    lang_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(1), Inches(0.4))
    tf = lang_box.text_frame
    p = tf.paragraphs[0]
    p.text = language
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['accent_blue']

    # Code
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(4.7))
    tf = code_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = COLORS['accent_green']

    return slide

def add_diagram_slide(prs, title, diagram_text, subtitle=None):
    """Add a slide with a text-based diagram."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_gray']

    # Diagram box
    y_start = Inches(1.4) if subtitle else Inches(1.1)
    diag_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y_start, Inches(9.4), Inches(5.0))
    diag_bg.fill.solid()
    diag_bg.fill.fore_color.rgb = COLORS['code_bg']
    diag_bg.line.color.rgb = COLORS['accent_purple']
    diag_bg.line.width = Pt(2)

    # Diagram text
    diag_box = slide.shapes.add_textbox(Inches(0.5), y_start + Inches(0.2), Inches(9), Inches(4.6))
    tf = diag_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    p.font.color.rgb = COLORS['text_white']

    return slide

def add_table_slide(prs, title, headers, rows, accent_color=None):
    """Add a slide with a table."""
    if accent_color is None:
        accent_color = COLORS['accent_blue']

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    # Calculate table dimensions
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = Inches(9)
    table_height = Inches(0.5) * num_rows
    col_width = table_width / num_cols

    # Add table
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.2), table_width, table_height).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['bg_dark']
        p.alignment = PP_ALIGN.CENTER

    # Style data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['bg_medium'] if row_idx % 2 == 0 else COLORS['bg_dark']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_white']
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    # Left column header
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = COLORS['accent_green']
    left_header.line.fill.background()
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['bg_dark']
    p.alignment = PP_ALIGN.CENTER

    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_white']
        p.space_before = Pt(6)

    # Right column header
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = COLORS['accent_red']
    right_header.line.fill.background()
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER

    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_white']
        p.space_before = Pt(6)

    return slide

def add_metrics_slide(prs, title, metrics):
    """Add a slide with big metrics."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['bg_dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']

    # Metrics in a row
    colors = [COLORS['accent_blue'], COLORS['accent_green'], COLORS['accent_orange'], COLORS['accent_purple']]
    num_metrics = min(len(metrics), 4)
    box_width = 9 / num_metrics

    for i, (value, label) in enumerate(metrics[:4]):
        x = Inches(0.5 + i * box_width)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(box_width - 0.2), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_medium']
        box.line.color.rgb = colors[i % len(colors)]
        box.line.width = Pt(3)

        # Value
        val_box = slide.shapes.add_textbox(x, Inches(1.8), Inches(box_width - 0.2), Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.alignment = PP_ALIGN.CENTER

        # Label
        lab_box = slide.shapes.add_textbox(x, Inches(3.0), Inches(box_width - 0.2), Inches(0.8))
        tf = lab_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_gray']
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Create the full AMO-Lean presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    add_title_slide(prs,
        "AMO-Lean",
        "Automatic Mathematical Optimizer in Lean\nCompilador de Optimización Formal para Primitivas Criptográficas")

    # ========== SLIDE 2: Agenda ==========
    add_content_slide(prs, "Agenda", [
        "- Qué es AMO-Lean (y qué NO es)",
        "- El DSL: MatExpr, VecExpr, SigmaExpr",
        "- Pipeline de Optimización",
        "- Ejemplos Concretos en ZK",
        "- Benchmarks y Eficiencia",
        "- Estado Actual vs Producción",
        "- Escenarios de Uso en zkVMs",
        "- Roadmap"
    ])

    # ========== SECTION 1: What is AMO-Lean ==========
    add_section_slide(prs, "Qué es AMO-Lean", 1)

    # Core concept
    add_diagram_slide(prs, "El Concepto Central", """
                              AMO-Lean Pipeline
    ╔═══════════════════════════════════════════════════════════════════╗
    ║                                                                   ║
    ║   Especificación       E-Graph              Código C/SIMD         ║
    ║   Matemática      →    Equality        →    Optimizado            ║
    ║   (DSL en Lean)        Saturation           (Correcto por         ║
    ║                        (Reglas               Construcción)        ║
    ║                         Verificadas)                              ║
    ║                                                                   ║
    ╚═══════════════════════════════════════════════════════════════════╝

    Similar a:
    • SPIRAL  → DSL para signal processing → C optimizado
    • Halide  → DSL para image processing → C/CUDA optimizado
    • CompCert → C → Assembly verificado
    """)

    # What it is NOT
    add_comparison_slide(prs, "AMO-Lean: Qué ES vs Qué NO ES",
        "ES", [
            "Optimizador para un DSL específico",
            "Generador de código C/SIMD",
            "Motor de reescritura verificado",
            "Herramienta para primitivas ZK",
            "E-Graph con equality saturation"
        ],
        "NO ES", [
            "Compilador de Lean arbitrario",
            "Una zkVM completa",
            "Una librería criptográfica",
            "Específico a un solo protocolo",
            "Magia que optimiza todo"
        ])

    # ========== SECTION 2: The DSL ==========
    add_section_slide(prs, "El DSL de AMO-Lean", 2)

    # DSL Types
    add_content_slide(prs, "Tipos del DSL", [
        "- Expr α — Expresiones escalares",
        "  - const, var, add, mul, pow",
        "- VecExpr α n — Expresiones vectoriales (longitud en el tipo)",
        "  - var, const, add, smul, concat, split",
        "- MatExpr α m n — Expresiones matriciales",
        "  - identity, compose, tensor (⊗), perm, elemwise",
        "- SigmaExpr — IR de loops (Σ-SPL de SPIRAL)",
        "  - gather, scatter, kernel, loop"
    ], COLORS['accent_purple'])

    # Code example
    add_code_slide(prs, "Ejemplo: FRI Fold en el DSL", """
-- Especificación FRI Fold como VecExpr
def foldSpec (n : Nat) : VecExpr F n :=
  let even  := VecExpr.var 0      -- Vector de entrada (even)
  let odd   := VecExpr.var 1      -- Vector de entrada (odd)
  let alpha := Expr.var 0         -- Escalar (challenge)

  -- fold(even, odd, α) = even + α * odd
  VecExpr.add even (VecExpr.smul alpha odd)

-- AMO-Lean genera automáticamente:
void fri_fold(const field_t* even, const field_t* odd,
              field_t* out, field_t alpha) {
    for (size_t i = 0; i < n; i++) {
        out[i] = field_add(even[i], field_mul(alpha, odd[i]));
    }
}
""", "lean → C")

    # MatExpr for FFT
    add_code_slide(prs, "MatExpr: FFT/NTT en O(log N) nodos", """
-- Fórmula de Cooley-Tukey (de SPIRAL):
-- DFT_{mn} = (DFT_m ⊗ I_n) · T^{mn}_n · (I_m ⊗ DFT_n) · L^{mn}_m

-- En AMO-Lean:
def ntt_spec (n : Nat) : MatExpr F n n :=
  let dft_half := MatExpr.dft (n/2)
  let twiddle  := MatExpr.twiddle n
  let stride   := MatExpr.stride n 2

  MatExpr.compose
    (MatExpr.tensor dft_half (MatExpr.identity 2))
    (MatExpr.compose twiddle
      (MatExpr.compose
        (MatExpr.tensor (MatExpr.identity 2) dft_half)
        stride))

-- Representa O(N log N) operaciones con O(log N) nodos AST
-- El E-Graph puede encontrar factorizaciones óptimas
""", "lean")

    # ========== SECTION 3: Optimization Pipeline ==========
    add_section_slide(prs, "Pipeline de Optimización", 3)

    # E-Graph explanation
    add_content_slide(prs, "E-Graph + Equality Saturation", [
        "- E-Graph: estructura que representa MUCHAS expresiones equivalentes",
        "- Equality Saturation: aplica TODAS las reglas hasta saturar",
        "- Después extrae la expresión de menor costo",
        "",
        "- Reglas verificadas formalmente (teoremas en Lean):",
        "  - Identity: x + 0 = x, x * 1 = x",
        "  - Zero propagation: x * 0 = 0",
        "  - Power rules: x^0 = 1, x^1 = x",
        "  - Factorization: a*b + a*c = a*(b+c)",
        "  - Constant folding: const a + const b = const (a+b)"
    ], COLORS['accent_green'])

    # Optimization results
    add_metrics_slide(prs, "Resultados de Optimización", [
        ("91.67%", "Reducción de\nOperaciones"),
        ("12+", "Reglas\nVerificadas"),
        ("0", "Sorry en\nTeoremas"),
        ("500/500", "Tests de\nEquivalencia")
    ])

    # ========== SECTION 4: Concrete Examples ==========
    add_section_slide(prs, "Ejemplos Concretos en ZK", 4)

    # FRI Fold
    add_content_slide(prs, "1. FRI Fold (Implementado)", [
        "- Operación core de FRI (Fast Reed-Solomon IOP)",
        "- fold(even, odd, α) = even + α * odd",
        "",
        "- Implementación escalar: field_goldilocks.h",
        "- Implementación AVX2: field_goldilocks_avx2.h",
        "- 4 elementos en paralelo (256-bit SIMD)",
        "",
        "- Usado en: Plonky2, Plonky3, STARK provers"
    ], COLORS['accent_blue'])

    # Poseidon2
    add_content_slide(prs, "2. Poseidon2 Hash (Implementado)", [
        "- Hash algebraico para ZK (Merkle trees, Fiat-Shamir)",
        "- S-Box: x^5 (BN254) o x^7 (Goldilocks)",
        "- MDS matrix mixing",
        "",
        "- Expresado como MatExpr con ElemOp:",
        "  - round = MDS · elemwise(sbox)",
        "",
        "- Constantes precomputadas para BN254",
        "- Usado en: STARKs, Polygon zkEVM, Scroll"
    ], COLORS['accent_orange'])

    # NTT
    add_content_slide(prs, "3. NTT/FFT (Parcialmente Implementado)", [
        "- Number Theoretic Transform para multiplicación de polinomios",
        "- Expresable como MatExpr con productos Kronecker",
        "",
        "- Estado actual:",
        "  - MatExpr con DFT, twiddles, permutaciones",
        "  - Sigma-SPL para lowering a loops",
        "  - CodeGen parcial",
        "",
        "- Usado en: Todos los provers polinomiales"
    ], COLORS['accent_purple'])

    # Polynomial evaluation
    add_code_slide(prs, "4. Evaluación de Polinomios (Expresable)", """
-- Horner's method: p(x) = a_n*x^n + ... + a_1*x + a_0
--                      = (...((a_n*x + a_{n-1})*x + ...)*x + a_0

def horner_spec (coeffs : Vec (Expr F) n) (x : Expr F) : Expr F :=
  coeffs.foldl
    (fun acc c => Expr.add c (Expr.mul acc x))
    (Expr.const 0)

-- El E-Graph puede:
-- 1. Detectar evaluaciones redundantes
-- 2. Factorizar subexpresiones comunes
-- 3. Aplicar optimizaciones específicas del campo

-- Usado en: FRI queries, constraint evaluation
""", "lean")

    # ========== SECTION 5: Benchmarks ==========
    add_section_slide(prs, "Benchmarks y Eficiencia", 5)

    # Main metrics
    add_metrics_slide(prs, "Speedups Medidos", [
        ("32.3x", "Lean → C\n(escalar)"),
        ("4.00x", "Scalar → AVX2\n(4-way SIMD)"),
        ("~129x", "Lean → C AVX2\n(combinado)"),
        ("568M", "elem/s\nGoldilocks")
    ])

    # Detailed benchmarks table
    add_table_slide(prs, "Benchmarks Detallados",
        ["Operación", "Sin Optimizar", "AMO-Lean", "Speedup"],
        [
            ["FRI Fold (Lean)", "28.4 ms", "0.88 ms", "32.3x"],
            ["FRI Fold (Scalar→AVX2)", "4.0 ms", "1.0 ms", "4.0x"],
            ["Goldilocks Mul", "~50 cycles", "~12 cycles", "4.2x"],
            ["E-Graph Reduction", "24 ops", "2 ops", "12x"],
        ])

    # Efficiency analysis
    add_content_slide(prs, "Análisis de Eficiencia", [
        "- AVX2 alcanza 100% del speedup teórico (4x para 4-way)",
        "",
        "- Comparación con código no optimizado formalmente:",
        "  - Código naive puede tener branches innecesarios",
        "  - AMO-Lean elimina dead code automáticamente",
        "  - Constant folding reduce operaciones en tiempo de compilación",
        "",
        "- Overhead de campo finito:",
        "  - Goldilocks vs UInt64: ~5x (aceptable)",
        "  - Goldilocks vs Barrett: ~2x mejor (reducción especializada)"
    ])

    # vs other implementations
    add_table_slide(prs, "Comparación: Goldilocks Implementations",
        ["Implementación", "Técnica", "Rendimiento Relativo"],
        [
            ["AMO-Lean AVX2", "Reducción especializada + SIMD", "100% (baseline)"],
            ["Plonky3 AVX2", "Similar (referencia)", "~100%"],
            ["Naive Barrett", "Reducción genérica", "~50%"],
            ["Naive sin SIMD", "Escalar simple", "~25%"],
        ],
        COLORS['accent_green'])

    # ========== SECTION 6: Current State ==========
    add_section_slide(prs, "Estado Actual vs Producción", 6)

    # Phase completion
    add_table_slide(prs, "Fases del Proyecto",
        ["Fase", "Descripción", "Estado", "Resultado"],
        [
            ["0", "Proof of Concept", "✅", "32.3x speedup"],
            ["1", "Goldilocks Field", "✅", "568 M elem/s"],
            ["2", "E-Graph Optimization", "✅", "91.67% reducción"],
            ["3", "AVX2 SIMD", "✅", "4.00x speedup"],
            ["4", "API Producción", "🔄", "En progreso"],
        ])

    # What works vs what's missing
    add_comparison_slide(prs, "Estado Actual",
        "FUNCIONA", [
            "FRI Fold completo (escalar + AVX2)",
            "E-Graph con 12+ reglas verificadas",
            "Poseidon2 spec + codegen",
            "Campo Goldilocks optimizado",
            "CI con 1000+ tests passing"
        ],
        "FALTA", [
            "NTT/FFT codegen completo",
            "Merkle tree optimizado",
            "API de usuario limpia",
            "Más campos (M31, BabyBear)",
            "Documentación de integración"
        ])

    # ========== SECTION 7: Use Cases ==========
    add_section_slide(prs, "Escenarios de Uso en zkVMs", 7)

    # zkVM integration diagram
    add_diagram_slide(prs, "Integración con zkVM", """
    ┌─────────────────────────────────────────────────────────────────────┐
    │                         zkVM Pipeline                                │
    ├─────────────────────────────────────────────────────────────────────┤
    │                                                                     │
    │  1. Execution    →  2. AIR           →  3. FRI Prover              │
    │     Trace            Constraints          (commitments)             │
    │     (witness)        (algebraic)                                    │
    │                                                                     │
    │                            ↓                                        │
    │                  ╔═════════════════════╗                            │
    │                  ║     AMO-Lean        ║                            │
    │                  ║  genera código      ║                            │
    │                  ║  optimizado para:   ║                            │
    │                  ║  • FRI fold         ║                            │
    │                  ║  • NTT/iNTT         ║                            │
    │                  ║  • Poseidon2        ║                            │
    │                  ║  • Merkle trees     ║                            │
    │                  ╚═════════════════════╝                            │
    │                            ↓                                        │
    │                  C/AVX2/AVX512 headers                              │
    │                                                                     │
    └─────────────────────────────────────────────────────────────────────┘
    """, "AMO-Lean genera los kernels de bajo nivel")

    # Use case 1
    add_content_slide(prs, "Caso de Uso 1: Prover Acelerado", [
        "- Problema: Provers ZK son lentos (minutos para programas grandes)",
        "- Bottlenecks: NTT, hashing, FRI queries",
        "",
        "- Solución con AMO-Lean:",
        "  - Especificar operaciones en DSL",
        "  - Generar kernels AVX2/AVX512 optimizados",
        "  - Integrar como headers .h",
        "",
        "- Beneficio: Speedup garantizado + corrección formal"
    ], COLORS['accent_blue'])

    # Use case 2
    add_content_slide(prs, "Caso de Uso 2: Multi-Campo", [
        "- Problema: Diferentes zkVMs usan diferentes campos",
        "  - Plonky2/3: Goldilocks (2^64 - 2^32 + 1)",
        "  - Polygon zkEVM: Goldilocks",
        "  - Scroll: BN254",
        "  - Risc0: M31 (Mersenne)",
        "",
        "- Solución con AMO-Lean:",
        "  - Misma especificación, diferentes campos",
        "  - #define FIELD_GOLDILOCKS / FIELD_BN254",
        "  - Código optimizado automáticamente para cada campo"
    ], COLORS['accent_green'])

    # Use case 3
    add_content_slide(prs, "Caso de Uso 3: Hardware Heterogéneo", [
        "- Problema: Diferentes CPUs tienen diferentes capacidades SIMD",
        "  - Intel Skylake: AVX2 (256-bit)",
        "  - Intel Ice Lake: AVX-512 (512-bit)",
        "  - AMD Zen4: AVX-512 + mejorado",
        "  - ARM: NEON / SVE",
        "",
        "- Solución (futuro) con AMO-Lean:",
        "  - Backend configurable por target",
        "  - Misma spec → diferentes implementaciones",
        "  - Runtime detection posible"
    ], COLORS['accent_orange'])

    # Code integration example
    add_code_slide(prs, "Ejemplo: Integración en Prover C", """
// Archivo generado por AMO-Lean
#include "fri_fold_goldilocks_avx2.h"
#include "ntt_goldilocks_avx2.h"
#include "poseidon2_goldilocks_avx2.h"

void prove_layer(ProverState* state) {
    // Operaciones optimizadas por AMO-Lean

    // 1. NTT forward
    ntt_forward_avx2(state->evaluations, state->n);

    // 2. Merkle commit con Poseidon2
    poseidon2_merkle_avx2(state->tree, state->evaluations, state->n);

    // 3. FRI fold
    fri_fold_avx2(state->folded,
                  state->even, state->odd,
                  state->alpha, state->n / 2);
}
""", "C")

    # ========== SECTION 8: Roadmap ==========
    add_section_slide(prs, "Roadmap", 8)

    # Future work
    add_content_slide(prs, "Próximos Pasos", [
        "- Phase 4: API de Producción",
        "  - Función compileToC(spec, config) limpia",
        "  - Documentación completa",
        "",
        "- Completar NTT/FFT CodeGen",
        "  - Radix-2 y Radix-4",
        "  - Optimización de twiddle factors",
        "",
        "- Más campos",
        "  - M31 (Mersenne 2^31 - 1)",
        "  - BabyBear (2^31 - 2^27 + 1)",
        "",
        "- AVX-512 backend"
    ])

    # Ideal API
    add_code_slide(prs, "API Ideal (Objetivo)", """
-- Usuario escribe especificación de alto nivel:
def my_prover_kernel : ProverSpec := {
  field := Goldilocks,
  operations := [
    .fri_fold { layers := 10 },
    .ntt { size := 2^20, radix := 4 },
    .poseidon2 { width := 3 }
  ],
  target := .AVX2
}

-- Un comando genera todo:
#generate my_prover_kernel "output/my_prover.h"

-- Resultado: headers .h listos para integrar
-- Con garantía de corrección formal
""", "lean (futuro)")

    # ========== FINAL SLIDES ==========

    # Summary metrics
    add_metrics_slide(prs, "Resumen: Lo que Logramos", [
        ("~129x", "Speedup\nTotal"),
        ("1061+", "Tests\nPassing"),
        ("4", "Fases\nCompletas"),
        ("0", "Sorry\nen Teoremas")
    ])

    # Key takeaways
    add_content_slide(prs, "Conclusiones", [
        "- AMO-Lean es un compilador de optimización formal",
        "  - NO código Lean arbitrario → DSL específico",
        "",
        "- Genera código C/SIMD con garantías de corrección",
        "  - Reglas de reescritura con teoremas formales",
        "  - E-Graph encuentra optimizaciones automáticamente",
        "",
        "- Aplicable a primitivas ZK:",
        "  - FRI fold, NTT, Poseidon2, evaluación de polinomios",
        "",
        "- Estado: 4/5 fases completas, falta API de producción"
    ])

    # Q&A
    add_title_slide(prs, "¿Preguntas?",
        "github.com/manuelpuebla/amo-lean")

    return prs

def main():
    """Generate and save the presentation."""
    # Ensure output directory exists
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "AMO-Lean_Technical_Talk.pptx")

    print("Generating AMO-Lean Technical Presentation...")
    prs = create_presentation()
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
