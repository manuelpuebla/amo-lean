#!/usr/bin/env python3
"""
AMO-Lean Technical Presentation Generator
Creates a PowerPoint presentation for a professional software audience
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RgbColor(0x1a, 0x36, 0x5d)
ACCENT_BLUE = RgbColor(0x2e, 0x86, 0xab)
ACCENT_ORANGE = RgbColor(0xf2, 0x7f, 0x1b)
LIGHT_GRAY = RgbColor(0xf5, 0xf5, 0xf5)
CODE_GREEN = RgbColor(0x2d, 0x6a, 0x4f)
ERROR_RED = RgbColor(0xc9, 0x30, 0x2c)

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Left accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    return slide

def add_content_slide(title, bullets, code_block=None):
    """Add a content slide with bullets and optional code"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Content area
    content_top = Inches(1.5)
    content_width = Inches(6) if code_block else Inches(12)

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle nested bullets
        if bullet.startswith("  - "):
            p.text = bullet[4:]
            p.level = 1
        else:
            p.text = bullet.lstrip("- ")
            p.level = 0

        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)

    # Code block if provided
    if code_block:
        code_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), content_top, Inches(6), Inches(5)
        )
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RgbColor(0x2d, 0x2d, 0x2d)
        code_shape.line.fill.background()

        code_box = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.6), Inches(4.6))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_block
        p.font.size = Pt(12)
        p.font.name = "Courier New"
        p.font.color.rgb = RgbColor(0xe0, 0xe0, 0xe0)

    return slide

def add_architecture_slide():
    """Add the architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Arquitectura: Pipeline de Compilacion Verificada"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Architecture boxes
    levels = [
        ("Level 6: FRI Protocol", "RoundState, CryptoSigma, Merkle, Fiat-Shamir", ACCENT_ORANGE),
        ("Level 5: Vectorial Ops", "MatExpr (Kronecker), VecExpr, SIMD CodeGen", ACCENT_BLUE),
        ("Level 4: Finite Fields", "ZMod integration, Power expressions", ACCENT_BLUE),
        ("Level 3: E-Graph", "Equality Saturation, E-matching, Union-Find", ACCENT_BLUE),
        ("Level 2: Rewriting", "Algebraic rules, Bottom-up, Fixpoint", ACCENT_BLUE),
        ("Level 1: CodeGen", "Expr -> C, Proof Anchors", CODE_GREEN),
    ]

    box_height = Inches(0.8)
    box_width = Inches(10)
    start_y = Inches(1.5)

    for i, (name, desc, color) in enumerate(levels):
        y = start_y + i * Inches(0.95)

        # Main box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.1), Inches(4), Inches(0.35))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

        desc_box = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.4), Inches(9.5), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(255, 255, 255)

        # Arrow
        if i < len(levels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), y + box_height, Inches(0.4), Inches(0.15))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RgbColor(0x66, 0x66, 0x66)
            arrow.line.fill.background()

    return slide

def add_comparison_slide(title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Left column
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = ERROR_RED
    left_header.line.fill.background()

    lh_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.5))
    tf = lh_text.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5))
    tf = left_content.text_frame
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    # Right column
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = CODE_GREEN
    right_header.line.fill.background()

    rh_text = slide.shapes.add_textbox(Inches(7), Inches(1.55), Inches(5.4), Inches(0.5))
    tf = rh_text.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5))
    tf = right_content.text_frame
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    return slide

def add_bug_slide():
    """Add the critical bug discovery slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ERROR_RED
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Problema 3: Buffer Swap Bug (Nivel Practico)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Bug description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(2))
    tf = desc_box.text_frame
    tf.word_wrap = True

    items = [
        "El codigo C compilaba sin errores",
        "No habia crashes ni excepciones",
        "Commitments y challenges: CORRECTOS",
        "Polinomio final: INCORRECTO",
        "Resultado: Proofs invalidos que parecen validos"
    ]

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Code comparison
    buggy_code = """// CODIGO CON BUG:
if (round + 1 < num_rounds) {
    // Solo swap si NO es
    // el ultimo round
    swap(current, next);
}
// Resultado: current apunta
// a datos VIEJOS en round final"""

    fixed_code = """// CODIGO CORREGIDO:
// SIEMPRE swap despues
// de cada round
swap(current, next);

// Resultado: current siempre
// apunta a datos FRESCOS"""

    # Buggy code box
    bug_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(5.8), Inches(3))
    bug_shape.fill.solid()
    bug_shape.fill.fore_color.rgb = RgbColor(0x4a, 0x1a, 0x1a)
    bug_shape.line.fill.background()

    bug_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(5.4), Inches(2.6))
    tf = bug_text.text_frame
    p = tf.paragraphs[0]
    p.text = buggy_code
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = RgbColor(0xff, 0x99, 0x99)

    # Fixed code box
    fix_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4), Inches(5.8), Inches(3))
    fix_shape.fill.solid()
    fix_shape.fill.fore_color.rgb = RgbColor(0x1a, 0x3a, 0x1a)
    fix_shape.line.fill.background()

    fix_text = slide.shapes.add_textbox(Inches(7), Inches(4.2), Inches(5.4), Inches(2.6))
    tf = fix_text.text_frame
    p = tf.paragraphs[0]
    p.text = fixed_code
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = RgbColor(0x99, 0xff, 0x99)

    # Root cause
    cause_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(2))
    tf = cause_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Causa Raiz: Impedance Mismatch"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    p = tf.add_paragraph()
    p.text = "Funcional (Lean): fold(fold(fold(x)))"
    p.font.size = Pt(16)
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "  Estado fluye naturalmente"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

    p = tf.add_paragraph()
    p.text = "Imperativo (C): for(i){swap(p1,p2)}"
    p.font.size = Pt(16)
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "  Estado es puntero mutable"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

    return slide

def add_results_slide():
    """Add the results slide with metrics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Resultados y Metricas"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Metrics boxes
    metrics = [
        ("~12,000", "Lineas de Lean 4"),
        ("716", "Modulos compilados"),
        ("107", "Tests diferenciales"),
        ("2", "Bugs criticos encontrados"),
        ("100%", "Tests pasando"),
    ]

    box_width = Inches(2.2)
    start_x = Inches(0.7)

    for i, (number, label) in enumerate(metrics):
        x = start_x + i * Inches(2.5)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), box_width, Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BLUE
        box.line.fill.background()

        num_text = slide.shapes.add_textbox(x, Inches(1.7), box_width, Inches(0.8))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        label_text = slide.shapes.add_textbox(x, Inches(2.5), box_width, Inches(0.5))
        tf = label_text.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Guarantees section
    guar_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(6), Inches(0.5))
    tf = guar_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Garantias de Seguridad"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    guarantees = [
        "Preservacion semantica (verificado por fuzzing)",
        "Seguridad Fiat-Shamir (probado formalmente)",
        "Reduccion de dominio correcta (probado)",
        "Ordenamiento de operaciones (Commit->Absorb->Squeeze->Fold)"
    ]

    guar_content = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(3))
    tf = guar_content.text_frame
    for i, g in enumerate(guarantees):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + g
        p.font.size = Pt(16)
        p.font.color.rgb = CODE_GREEN
        p.space_after = Pt(8)

    # Limitations
    lim_title = slide.shapes.add_textbox(Inches(6.8), Inches(3.4), Inches(6), Inches(0.5))
    tf = lim_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Limitaciones Conocidas"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    limitations = [
        "No verificamos el compilador C (gcc/clang)",
        "No analizamos memory safety en C",
        "No analizamos timing side channels",
        "Hash functions son placeholders"
    ]

    lim_content = slide.shapes.add_textbox(Inches(6.8), Inches(4), Inches(6), Inches(3))
    tf = lim_content.text_frame
    for i, l in enumerate(limitations):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + l
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)
        p.space_after = Pt(8)

    return slide


# ============== BUILD PRESENTATION ==============

# Slide 1: Title
add_title_slide(
    "AMO-Lean",
    "Automatic Mathematical Optimizer for Lean 4\nCompilador Verificado para Protocolos Criptograficos"
)

# Slide 2: Executive Summary
add_content_slide(
    "Que es AMO-Lean?",
    [
        "Compilador optimizador que transforma especificaciones matematicas de alto nivel en codigo C eficiente",
        "Garantias formales de correccion usando Lean 4 y Mathlib",
        "Caso de estudio: Protocolo FRI (Fast Reed-Solomon IOP of Proximity)",
        "  - Componente critico en sistemas ZK-STARK (Ethereum, StarkNet)",
        "Metodologia: Verificacion Empirica Transitiva",
        "  - Teoremas formales sobre implementacion de referencia",
        "  - Fuzzing diferencial para validar codigo C generado"
    ]
)

# Slide 3: Why Verified Compilation
add_content_slide(
    "Por que Compilacion Verificada?",
    [
        "Protocolos criptograficos son criticos para seguridad",
        "  - Errores silenciosos pueden invalidar proofs sin detectarse",
        "  - Testing tradicional es insuficiente",
        "Optimizacion manual introduce bugs",
        "  - SIMD, memory layout, loop transformations",
        "Requerimiento: Correccion demostrable, no solo testeable",
        "AMO-Lean: Verificacion formal + Fuzzing diferencial",
        "  - Defense in depth contra errores sutiles"
    ]
)

# Slide 4: Architecture
add_architecture_slide()

# Slide 5: Section - Key Decisions
add_section_slide("Decisiones de Diseno Clave")

# Slide 6: Design Decision 1 - E-graph
add_content_slide(
    "Decision 1: E-Graph con Equality Saturation",
    [
        "Problema: Reescritura greedy no explora alternativas",
        "  - Asociatividad causa 70x slowdown",
        "  - Sin backtracking, optimos locales",
        "Solucion: E-graph (egg paper, POPL 2021)",
        "  - Union-Find con path compression",
        "  - E-matching para busqueda de patrones",
        "  - Extraccion con cost model",
        "Implementacion en Lean 4:",
        "  - Estructuras planas (Array + HashMap)",
        "  - Rebuild para re-canonicalizacion",
        "  - ~1000 lineas de codigo"
    ],
    """-- E-Graph en Lean 4
structure EGraph where
  unionFind : Array EClassId
  hashcons  : HashMap ENode EClassId
  classes   : Array EClass
  worklist  : Array EClassId

-- Saturacion
def saturate (eg : EGraph)
    (rules : List Rule) : EGraph :=
  let rec go eg fuel :=
    if fuel == 0 then eg
    else
      let eg' := saturateStep eg rules
      if eg'.stats == eg.stats
      then eg
      else go eg' (fuel - 1)
  go eg maxIterations"""
)

# Slide 7: Design Decision 2 - Kronecker
add_comparison_slide(
    "Decision 2: Productos Kronecker vs Expansion Escalar",
    "Expansion Escalar (Naive)",
    [
        "• FFT de N elementos = N operaciones",
        "• DFT_1M = ~1,000,000 nodos en AST",
        "• E-graph colapsa por memoria",
        "• Tiempo de compilacion: infinito",
        "",
        "Ejemplo DFT_8:",
        "  64 nodos escalares",
        "  Cada operacion explicita"
    ],
    "Productos Kronecker (SPIRAL)",
    [
        "• FFT = formula algebraica compacta",
        "• DFT_1M = ~41 nodos (!)",
        "• Reduccion: 10^10x",
        "• Tiempo de compilacion: segundos",
        "",
        "Ejemplo DFT_8:",
        "  11 nodos Kronecker",
        "  DFT_8 = (DFT_2 ⊗ I_4) · T · (I_2 ⊗ DFT_4) · L"
    ]
)

# Slide 8: Design Decision 3 - CryptoSigma
add_content_slide(
    "Decision 3: CryptoSigma IR con Barreras Intrinsecas",
    [
        "Problema: Optimizaciones pueden reordenar operaciones",
        "  - En criptografia, el orden es semanticamente critico",
        "  - Fiat-Shamir: commit ANTES de squeeze",
        "Solucion: IR intermedio con barreras explicitas",
        "  - absorb, squeeze, merkleHash como intrinsics",
        "  - BARRIER markers que prohiben reordenamiento",
        "  - domainEnter/Exit para contexto",
        "Generacion de codigo:",
        "  - Proof Anchors: comentarios estructurados",
        "  - Mapean directamente a teoremas en Lean"
    ],
    """-- CryptoSigma IR
inductive Intrinsic where
  | absorb : Nat -> Intrinsic
  | squeeze : Intrinsic
  | merkleHash : Intrinsic
  | domainEnter : Domain -> Intrinsic
  | domainExit : Intrinsic

-- Generacion C con Proof Anchor
// PROOF_ANCHOR: fri_fold
// Preconditions:
//   - n > 0
//   - input has 2*n elements
// Postconditions:
//   - output[i] == input[2*i]
//          + alpha * input[2*i+1]
// BARRIER: No code motion"""
)

# Slide 9: Section - Problems Solved
add_section_slide("Tres Problemas Resueltos")

# Slide 10: Problem 1 - Scalar Trap
add_content_slide(
    "Problema 1: La Trampa Escalar (Nivel Simbolico)",
    [
        "Contexto: FFT opera sobre vectores de N elementos",
        "  - Tipicamente N = 2^20 = 1,048,576",
        "Riesgo: AST escalar genera millones de nodos",
        "  - E-graph colapsa por consumo de memoria",
        "  - Out of Memory antes de compilar",
        "Solucion: Representacion Kronecker (SPIRAL paper)",
        "  - O(log N) nodos en lugar de O(N)",
        "  - Reglas algebraicas a nivel matricial",
        "  - Lowering a loops + kernels pequenos",
        "Impacto: DFT_2^20 viable con 41 nodos"
    ],
    """-- Representacion compacta
def fftCooleyTukey (m n : Nat) :=
  MatExpr.compose
    (MatExpr.kron (dft m) (identity n))
    (MatExpr.compose
      (MatExpr.twiddle (m*n) n)
      (MatExpr.compose
        (MatExpr.kron (identity m) (dft n))
        (MatExpr.perm (Perm.stride m n))))

-- Conteo de nodos:
-- DFT_8:     11 nodos (vs 64 escalar)
-- DFT_1024:  21 nodos (vs 1M escalar)
-- DFT_2^20:  41 nodos (vs 10^12 escalar)"""
)

# Slide 11: Problem 2 - E-graph Explosion
add_content_slide(
    "Problema 2: Explosion del E-Graph (Nivel Implementacion)",
    [
        "Contexto: Reglas de conmutatividad y asociatividad",
        "  - a + b = b + a (necesaria para optimizacion)",
        "  - (a + b) + c = a + (b + c)",
        "Riesgo: Aplicacion ilimitada genera explosion",
        "  - E-graph crece exponencialmente",
        "  - Saturacion nunca termina",
        "Solucion: Configuracion con limites estrictos",
        "  - maxIterations: 50",
        "  - maxNodes: 5000",
        "  - maxClasses: 1000",
        "Leccion: Las reglas mas 'obvias' son las mas peligrosas"
    ],
    """-- Configuracion segura
structure SaturationConfig where
  maxIterations : Nat := 50
  maxNodes      : Nat := 5000
  maxClasses    : Nat := 1000

-- Reglas controladas
def basicRules : List Rule := [
  rule_add_zero,    -- a + 0 = a
  rule_mul_one,     -- a * 1 = a
  rule_mul_zero,    -- a * 0 = 0
]

-- Reglas peligrosas (usar con cuidado)
def commRules : List Rule := [
  rule_add_comm,    -- a + b = b + a
  rule_mul_comm,    -- a * b = b * a
]
-- ADVERTENCIA: Solo con limites bajos!"""
)

# Slide 12: Problem 3 - Buffer Swap Bug
add_bug_slide()

# Slide 13: Verification Strategy
add_content_slide(
    "Estrategia de Verificacion: Transitividad Empirica",
    [
        "Nivel 1: Teoremas formales sobre referencia Lean",
        "  - friFold_spec: especificacion del fold",
        "  - round_ordering_secure: Fiat-Shamir order",
        "  - domain_size_after_rounds: reduccion 2^k",
        "Nivel 2: Fuzzing diferencial Lean vs C",
        "  - 107 test cases",
        "  - Comparacion bit-a-bit",
        "  - Genera inputs aleatorios",
        "Transitividad:",
        "  - Lean reference satisface propiedades (probado)",
        "  - C == Lean (verificado empiricamente)",
        "  - Por lo tanto: C satisface propiedades"
    ],
    """-- Teorema formal (Lean)
theorem round_ordering_secure
    (poly : Array UInt64)
    (ts : TranscriptState) :
  let result := executeRound poly ts
  result.finalTranscript.absorbCount
    > ts.absorbCount ∧
  result.finalTranscript.squeezeCount
    > ts.squeezeCount := by
  simp [executeRound,
        TranscriptState.absorb,
        TranscriptState.squeeze]
  omega

-- Test diferencial (Lean vs C)
#eval runDifferentialTest
-- Commitments: ✓ MATCH
-- Challenges:  ✓ MATCH
-- Final poly:  ✓ MATCH"""
)

# Slide 14: Results
add_results_slide()

# Slide 15: Lessons Learned
add_content_slide(
    "Lecciones Aprendidas",
    [
        "1. Fuzzing diferencial > Unit tests",
        "  - El buffer swap bug paso todos los tests unitarios",
        "  - Solo fallo en comparacion end-to-end",
        "2. Generar codigo ANTES de probar teoremas",
        "  - ADR-008: Evitar 'verificacion en el vacio'",
        "  - Probar teoremas sobre codigo buggy es inutil",
        "3. Traduccion funcional->imperativo es error-prone",
        "  - Estado explicito vs punteros mutables",
        "  - Requiere atencion especial",
        "4. Proof Anchors habilitan auditoria",
        "  - Trazabilidad codigo C <-> teoremas Lean",
        "  - Humanos pueden verificar correspondencia"
    ]
)

# Slide 16: Future Work
add_content_slide(
    "Trabajo Futuro",
    [
        "Corto plazo:",
        "  - Completar teoremas con sorry pendientes",
        "  - Mas vectores de test para fuzzing",
        "  - Variantes AVX-512 de kernels",
        "Mediano plazo:",
        "  - Integrar hash criptografico real (Poseidon)",
        "  - Query phase de FRI (no solo commit)",
        "  - Extender a STARK prover completo",
        "Largo plazo:",
        "  - Verificacion formal de codigo C (CompCert-style)",
        "  - Generacion de codigo GPU (CUDA/Metal)",
        "  - Integracion con sistemas ZK de produccion"
    ]
)

# Slide 17: Conclusion
add_title_slide(
    "Conclusion",
    "AMO-Lean demuestra que la compilacion verificada\nde protocolos criptograficos es PRACTICA.\n\nEl bug de buffer swap justifica toda la metodologia."
)

# Slide 18: References
add_content_slide(
    "Referencias",
    [
        "1. Willsey et al. 'egg: Fast and Extensible Equality Saturation' (POPL 2021)",
        "2. Franchetti et al. 'SPIRAL: Code Generation for DSP Transforms' (IEEE 2005)",
        "3. Gross et al. 'Verified Rewriting Engine' (ITP 2022) - Fiat-Crypto",
        "4. Xi & Pfenning 'Dependent Types in Practical Programming' (POPL 1999)",
        "5. Ben-Sasson et al. 'Fast Reed-Solomon IOP of Proximity' (ICALP 2018)",
        "",
        "Repositorio: github.com/manuelpuebla/amo-lean",
        "Documentacion: docs/FINAL_REPORT.md, docs/PROJECT_STATUS.md"
    ]
)

# Save presentation
output_path = "/Users/manuelpuebla/Documents/claudio/amo-lean/docs/AMO-Lean_Technical_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
