#!/usr/bin/env python3
"""
Update AMO-Lean presentation with zkVM-focused future work
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Load existing presentation
prs = Presentation("/Users/manuelpuebla/Documents/claudio/amo-lean/docs/AMO-Lean_Technical_Presentation.pptx")

# Color scheme
DARK_BLUE = RgbColor(0x1a, 0x36, 0x5d)
ACCENT_BLUE = RgbColor(0x2e, 0x86, 0xab)
ACCENT_ORANGE = RgbColor(0xf2, 0x7f, 0x1b)
CODE_GREEN = RgbColor(0x2d, 0x6a, 0x4f)
ERROR_RED = RgbColor(0xc9, 0x30, 0x2c)
CRITICAL_PURPLE = RgbColor(0x7b, 0x2d, 0x8e)

def add_content_slide_after(prs, index, title, bullets, code_block=None):
    """Add a content slide at specific position"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Move slide to correct position
    # (slides are added at the end, we need to reorder)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Content area
    content_top = Inches(1.5)
    content_width = Inches(6) if code_block else Inches(12)

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), content_top, content_width, Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle nested bullets and formatting
        if bullet.startswith("  - "):
            p.text = bullet[4:]
            p.level = 1
        elif bullet.startswith("    - "):
            p.text = bullet[6:]
            p.level = 2
        else:
            p.text = bullet.lstrip("- ")
            p.level = 0

        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)

    # Code block if provided
    if code_block:
        code_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), content_top, Inches(6), Inches(5)
        )
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RgbColor(0x2d, 0x2d, 0x2d)
        code_shape.line.fill.background()

        code_box = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.6), Inches(4.6))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_block
        p.font.size = Pt(11)
        p.font.name = "Courier New"
        p.font.color.rgb = RgbColor(0xe0, 0xe0, 0xe0)

    return slide

# Delete old "Trabajo Futuro" slide (slide 16, index 15)
# and replace with new slides

# First, let's find and update slide 16 (Future Work)
slide_16 = prs.slides[15]

# Clear existing shapes except header
shapes_to_remove = []
for shape in slide_16.shapes:
    shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Rebuild slide 16 with new content
# Header bar
header = slide_16.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = slide_16.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Trabajo Futuro: Ruta hacia zkVM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)

# Priority boxes
priorities = [
    ("#1 CRÍTICO", "Poseidon/Rescue Hash", "6-10 semanas", "Habilita recursion eficiente en zkVM", CRITICAL_PURPLE),
    ("#2 ALTO", "Backend CUDA/GPU", "3-6 meses", "Proof generation en granjas GPU", ERROR_RED),
    ("#3 MEDIO", "Variantes AVX-512", "2-3 semanas", "Optimizacion incremental CPU", ACCENT_ORANGE),
    ("#4 BAJO", "FRI Query Phase", "4-6 semanas", "Completa el protocolo", ACCENT_BLUE),
]

start_y = Inches(1.5)
box_height = Inches(1.3)

for i, (priority, name, time, desc, color) in enumerate(priorities):
    y = start_y + i * Inches(1.45)

    # Priority badge
    badge = slide_16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(1.8), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()

    badge_text = slide_16.shapes.add_textbox(Inches(0.5), y + Inches(0.08), Inches(1.8), Inches(0.4))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = priority
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Name
    name_text = slide_16.shapes.add_textbox(Inches(2.5), y, Inches(4), Inches(0.5))
    tf = name_text.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Time estimate
    time_text = slide_16.shapes.add_textbox(Inches(2.5), y + Inches(0.5), Inches(2), Inches(0.4))
    tf = time_text.text_frame
    p = tf.paragraphs[0]
    p.text = "⏱ " + time
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

    # Description
    desc_text = slide_16.shapes.add_textbox(Inches(2.5), y + Inches(0.85), Inches(10), Inches(0.4))
    tf = desc_text.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(0x44, 0x44, 0x44)

# Add note about roadmap
note_box = slide_16.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
tf = note_box.text_frame
p = tf.paragraphs[0]
p.text = "Ver: docs/ZKVM_ROADMAP.pdf para planificacion detallada por fases"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

# Now add a new slide about Poseidon Challenge (insert after slide 16)
# We'll add it at the end and note its position

slide_layout = prs.slide_layouts[6]
poseidon_slide = prs.slides.add_slide(slide_layout)

# Header
header = poseidon_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header.fill.solid()
header.fill.fore_color.rgb = CRITICAL_PURPLE
header.line.fill.background()

title_box = poseidon_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "El Desafio Poseidon: Por que es Critico"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)

# Left column - The problem
problem_title = poseidon_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
tf = problem_title.text_frame
p = tf.paragraphs[0]
p.text = "El Problema"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

problem_content = poseidon_slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(2.5))
tf = problem_content.text_frame
tf.word_wrap = True

problems = [
    "Poseidon = Lineal + No-Lineal",
    "  Capas MDS (matrices) ✓ Kronecker OK",
    "  S-boxes (x^5) ✗ Rompe linearidad",
    "zkVMs necesitan recursion eficiente",
    "  Una prueba verifica otra prueba",
    "  Hash es el cuello de botella",
    "Sin Poseidon optimizado = zkVM lenta"
]

for i, item in enumerate(problems):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    if item.startswith("  "):
        p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

# Right column - Our approach
approach_title = poseidon_slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5))
tf = approach_title.text_frame
p = tf.paragraphs[0]
p.text = "Nuestra Solucion"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# Code block
code_shape = poseidon_slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(2.1), Inches(5.8), Inches(3.5)
)
code_shape.fill.solid()
code_shape.fill.fore_color.rgb = RgbColor(0x2d, 0x2d, 0x2d)
code_shape.line.fill.background()

code_text = """-- Extender MatExpr
inductive MatExpr where
  | ...existentes...
  | elemwise : (Expr -> Expr)
             -> MatExpr n n  -- NUEVO

-- Poseidon round
def poseidonRound (state : MatExpr t t) :=
  MatExpr.compose
    (MatExpr.elemwise (x => x^5)) -- S-box
    (MatExpr.compose
      mdsMatrix                    -- Mezcla
      (MatExpr.add state consts)) -- RC"""

code_box = poseidon_slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.4), Inches(3.1))
tf = code_box.text_frame
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(12)
p.font.name = "Courier New"
p.font.color.rgb = RgbColor(0xe0, 0xe0, 0xe0)

# Impact box at bottom
impact_shape = poseidon_slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2)
)
impact_shape.fill.solid()
impact_shape.fill.fore_color.rgb = RgbColor(0xf0, 0xf0, 0xf0)
impact_shape.line.fill.background()

impact_text = poseidon_slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(1))
tf = impact_text.text_frame
p = tf.paragraphs[0]
p.text = "Impacto en zkVM: Con Poseidon optimizado, la recursion de pruebas pasa de minutos a segundos."
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p = tf.add_paragraph()
p.text = "Esto habilita: verificacion en tiempo real, pruebas incrementales, rollups mas eficientes."
p.font.size = Pt(16)
p.font.color.rgb = RgbColor(0x44, 0x44, 0x44)

# Add zkVM ecosystem slide
zkvm_slide = prs.slides.add_slide(slide_layout)

# Header
header = zkvm_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header.fill.solid()
header.fill.fore_color.rgb = DARK_BLUE
header.line.fill.background()

title_box = zkvm_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "AMO-Lean en el Ecosistema zkVM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 255, 255)

# Description
desc_box = zkvm_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
tf = desc_box.text_frame
p = tf.paragraphs[0]
p.text = "AMO-Lean genera los primitivos criptograficos verificados que las zkVMs necesitan"
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(0x44, 0x44, 0x44)

# zkVM boxes
zkvms = [
    ("SP1 (Succinct)", "RISC-V zkVM, recursion via Poseidon"),
    ("RISC Zero", "RISC-V zkVM, pruebas STARK"),
    ("Valida", "zkVM optimizada para EVM"),
    ("Cairo / Stone", "StarkNet prover, Poseidon-heavy"),
]

start_y = Inches(2.5)
for i, (name, desc) in enumerate(zkvms):
    y = start_y + i * Inches(1.1)

    box = zkvm_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(5.5), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BLUE
    box.line.fill.background()

    name_text = zkvm_slide.shapes.add_textbox(Inches(0.7), y + Inches(0.1), Inches(5), Inches(0.4))
    tf = name_text.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    desc_text = zkvm_slide.shapes.add_textbox(Inches(0.7), y + Inches(0.5), Inches(5), Inches(0.4))
    tf = desc_text.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(0xdd, 0xdd, 0xdd)

# What AMO-Lean provides
provides_title = zkvm_slide.shapes.add_textbox(Inches(6.5), Inches(2.5), Inches(6), Inches(0.5))
tf = provides_title.text_frame
p = tf.paragraphs[0]
p.text = "AMO-Lean Provee:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

provides = [
    "✓ FRI fold verificado (STARK core)",
    "✓ Merkle trees optimizados",
    "✓ Fiat-Shamir seguro (probado)",
    "→ Poseidon verificado (proximo)",
    "→ GPU kernels verificados (futuro)",
]

provides_box = zkvm_slide.shapes.add_textbox(Inches(6.5), Inches(3.1), Inches(6), Inches(3))
tf = provides_box.text_frame
for i, item in enumerate(provides):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(18)
    p.space_after = Pt(12)
    if item.startswith("✓"):
        p.font.color.rgb = CODE_GREEN
    else:
        p.font.color.rgb = ACCENT_ORANGE

# Value prop
value_box = zkvm_slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.7))
tf = value_box.text_frame
p = tf.paragraphs[0]
p.text = "Propuesta de valor: Primitivos criptograficos con garantias formales, listos para integracion."
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# Save updated presentation
prs.save("/Users/manuelpuebla/Documents/claudio/amo-lean/docs/AMO-Lean_Technical_Presentation.pptx")
print("Presentation updated successfully!")
